import requests
from bs4 import BeautifulSoup
from io import BytesIO
import pdfplumber
import docx
import pandas as pd


def extract_html(url):
    response = requests.get(url, timeout=10)
    response.raise_for_status()

    soup = BeautifulSoup(response.text, "html.parser")

    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()

    return soup.get_text(separator=" ", strip=True)


def extract_pdf(content):
    text = ""
    with pdfplumber.open(BytesIO(content)) as pdf:
        for page in pdf.pages:
            text += page.extract_text() or ""
    return text


def extract_docx(content):
    doc = docx.Document(BytesIO(content))
    return "\n".join([p.text for p in doc.paragraphs])


def extract_xlsx(content):
    df = pd.read_excel(BytesIO(content))
    return df.to_string()


def extract_pptx(content):
    from pptx import Presentation

    prs = Presentation(BytesIO(content))
    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text
